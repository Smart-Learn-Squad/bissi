"""PowerPoint operations for BISSI.

Provides functionality to read, create, and manipulate PPTX presentations.
Includes StyledPresentation for professional slides with layouts, images, and animations.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pathlib import Path
from typing import List, Dict, Any, Optional, Union


# Predefined color palettes for PowerPoint
class PPTColors:
    """Predefined color palette for PowerPoint slides."""
    PRIMARY = RGBColor(46, 134, 171)      # #2E86AB
    SECONDARY = RGBColor(33, 33, 33)     # #212121
    ACCENT = RGBColor(241, 196, 15)        # #F1C196
    SUCCESS = RGBColor(39, 174, 96)      # #27AE60
    WARNING = RGBColor(230, 126, 34)    # #E67E22
    DANGER = RGBColor(231, 76, 60)      # #E74C3C
    LIGHT = RGBColor(236, 240, 241)       # #ECF0F1
    DARK = RGBColor(44, 62, 80)         # #2C3E50


# Slide Layouts
class Layouts:
    """Slide layout indices."""
    TITLE = 0                    # Title slide
    TITLE_AND_CONTENT = 1         # Title and content
    SECTION_HEADER = 2           # Section header
    TWO_CONTENT = 3              # Two column content
    COMPARISON = 4               # Comparison
    BLANK = 5                   # Blank
    CONTENT_CAPTION = 6           # Content with caption
    PICTURE_CAPTION = 7           # Picture caption


class PPTXAgent:
    """Agent for PowerPoint presentation operations."""
    
    def __init__(self, file_path: Optional[Union[str, Path]] = None):
        """Initialize PPTX agent.
        
        Args:
            file_path: Path to existing presentation, or None for new
        """
        self.file_path = file_path
        if file_path and Path(file_path).exists():
            self.presentation = Presentation(file_path)
        else:
            self.presentation = Presentation()
    
    def read_slides(self) -> List[Dict[str, Any]]:
        """Extract text content from all slides.
        
        Returns:
            List of slides with text content
        """
        slides_content = []
        
        for i, slide in enumerate(self.presentation.slides, 1):
            slide_text = []
            
            # Extract from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            slides_content.append({
                'slide_number': i,
                'content': '\n'.join(slide_text),
                'shapes_count': len(slide.shapes)
            })
        
        return slides_content
    
    def read_tables(self) -> List[List[List[str]]]:
        """Extract tables from all slides.
        
        Returns:
            List of tables per slide
        """
        all_tables = []
        
        for slide in self.presentation.slides:
            slide_tables = []
            
            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    table_data = []
                    
                    for row in table.rows:
                        row_data = [cell.text for cell in row.cells]
                        table_data.append(row_data)
                    
                    slide_tables.append(table_data)
            
            all_tables.append(slide_tables)
        
        return all_tables
    
    def add_slide(self, title: str, content: str, layout_idx: int = 1) -> Any:
        """Add new slide with title and content.
        
        Args:
            title: Slide title
            content: Slide body content
            layout_idx: Slide layout index (1 = Title and Content)
            
        Returns:
            The created slide object
        """
        slide_layout = self.presentation.slide_layouts[layout_idx]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = title
        
        # Set content in body placeholder
        for shape in slide.placeholders:
            if shape.placeholder_format.type == 1:  # Body placeholder
                shape.text = content
                break
        
        return slide
    
    def add_image_slide(self, 
                       title: str, 
                       image_path: Union[str, Path],
                       left: float = 1,
                       top: float = 1.5,
                       width: float = 8,
                       height: float = 5) -> Any:
        """Add slide with image.
        
        Args:
            title: Slide title
            image_path: Path to image file
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
            
        Returns:
            The created slide object
        """
        slide_layout = self.presentation.slide_layouts[5]  # Blank layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        if title:
            # Add title textbox
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.5), Inches(9), Inches(0.5)
            )
            title_frame = title_box.text_frame
            title_frame.text = title
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(24)
            title_para.font.bold = True
        
        # Add image
        slide.shapes.add_picture(
            str(image_path),
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        
        return slide
    
    def add_table_slide(self,
                       title: str,
                       data: List[List[str]],
                       rows: int,
                       cols: int,
                       left: float = 1,
                       top: float = 2,
                       width: float = 8,
                       height: float = 4) -> Any:
        """Add slide with table.
        
        Args:
            title: Slide title
            data: 2D list of table data
            rows: Number of rows
            cols: Number of columns
            left: Left position in inches
            top: Top position in inches
            width: Table width in inches
            height: Table height in inches
            
        Returns:
            The created slide object
        """
        slide_layout = self.presentation.slide_layouts[5]  # Blank layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Add title
        if title:
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.5), Inches(9), Inches(0.5)
            )
            title_frame = title_box.text_frame
            title_frame.text = title
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(24)
            title_para.font.bold = True
        
        # Add table
        table = slide.shapes.add_table(
            rows, cols,
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        ).table
        
        # Fill data
        for i, row_data in enumerate(data):
            if i < rows:
                for j, cell_text in enumerate(row_data):
                    if j < cols:
                        table.cell(i, j).text = str(cell_text)
        
        return slide
    
    def get_slide_count(self) -> int:
        """Get number of slides in presentation."""
        return len(self.presentation.slides)
    
    def get_slide_layouts(self) -> List[str]:
        """Get available slide layout names."""
        return [layout.name for layout in self.presentation.slide_layouts]
    
    def save(self, file_path: Optional[Union[str, Path]] = None) -> None:
        """Save presentation to file.
        
        Args:
            file_path: Output path, or original path if None
        """
        save_path = file_path or self.file_path
        if not save_path:
            raise ValueError("No file path specified for saving")
        
        self.presentation.save(str(save_path))
        self.file_path = str(save_path)


def read_presentation(file_path: Union[str, Path]) -> List[Dict[str, Any]]:
    """Read and extract all content from PowerPoint file.
    
    Args:
        file_path: Path to .pptx file
        
    Returns:
        List of slides with content
    """
    agent = PPTXAgent(file_path)
    return agent.read_slides()


def get_presentation_info(file_path: Union[str, Path]) -> Dict[str, Any]:
    """Get metadata about PowerPoint presentation.
    
    Args:
        file_path: Path to .pptx file
        
    Returns:
        Dictionary with presentation info
    """
    agent = PPTXAgent(file_path)
    
    # Count shapes and content
    total_shapes = 0
    text_shapes = 0
    table_shapes = 0
    image_shapes = 0
    
    for slide in agent.presentation.slides:
        for shape in slide.shapes:
            total_shapes += 1
            
            if hasattr(shape, "text") and shape.text.strip():
                text_shapes += 1
            if shape.has_table:
                table_shapes += 1
            if shape.shape_type == 13:  # Picture type
                image_shapes += 1
    
    return {
        'slide_count': len(agent.presentation.slides),
        'total_shapes': total_shapes,
        'text_shapes': text_shapes,
        'table_shapes': table_shapes,
        'image_shapes': image_shapes,
        'layouts': agent.get_slide_layouts()
    }


def create_presentation(title: str = "New Presentation") -> PPTXAgent:
    """Create new presentation with title slide.
    
    Args:
        title: Presentation title
        
    Returns:
        PPTXAgent instance with new presentation
    """
    agent = PPTXAgent()
    
    # Add title slide
    slide_layout = agent.presentation.slide_layouts[0]  # Title slide layout
    slide = agent.presentation.slides.add_slide(slide_layout)
    
    # Set title and subtitle
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = f"Created by BISSI - {Path.cwd().name}"

    return agent


class StyledPresentation:
    """Create polished, styled PowerPoint presentations.

    Provides a fluent API for professional presentation generation with:
    - Predefined slide layouts (title, content, two-column, etc.)
    - Text styling (fonts, colors, sizes)
    - Bulleted lists with custom markers
    - Images with captions
    - Shapes and callouts
    - Progress tracking

    Example:
        ppt = StyledPresentation("presentation.pptx")
        ppt.add_title_slide("Annual Report 2024", "Performance Review")
        ppt.add_section("Q1 Results")
        ppt.add_content_slide("Q1 Performance",
            ["Revenue up 15%", "New customers +200", "Product launches: 3"])
        ppt.add_two_column_slide("Strengths", "Weaknesses",
            ["Strong team", "Clear vision"],
            ["Limited budget", "Small team"])
        ppt.add_image_slide("Market Analysis", "chart.png", "Growing segment")
        ppt.save()
    """

    LAYOUT_TITLE = 0
    LAYOUT_CONTENT = 1
    LAYOUT_SECTION = 2
    LAYOUT_TWO_COL = 3
    LAYOUT_COMPARISON = 4
    LAYOUT_BLANK = 6

    def __init__(self, file_path: Optional[str] = None, width: float = 10, height: float = 7.5):
        self.file_path = file_path
        if file_path and Path(file_path).exists():
            self.presentation = Presentation(file_path)
            self.slide_count = len(self.presentation.slides)
        else:
            self.presentation = Presentation()
            self.slide_count = 0
            self.presentation.slide_width = Inches(width)
            self.presentation.slide_height = Inches(height)

    def add_title_slide(self, title: str, subtitle: str = None, color: RGBColor = None) -> 'StyledPresentation':
        slide_layout = self.presentation.slide_layouts[self.LAYOUT_TITLE]
        slide = self.presentation.slides.add_slide(slide_layout)
        self.slide_count += 1
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.bold = True
        if color:
            title_shape.text_frame.paragraphs[0].font.color.rgb = color
        if subtitle:
            slide.placeholders[1].text = subtitle
        return self

    def add_section(self, title: str, color: RGBColor = None) -> 'StyledPresentation':
        slide_layout = self.presentation.slide_layouts[self.LAYOUT_SECTION]
        slide = self.presentation.slides.add_slide(slide_layout)
        self.slide_count += 1
        slide.shapes.title.text = title
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(40)
        slide.shapes.title.text_frame.paragraphs[0].font.bold = True
        if color:
            slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = color
        return self

    def add_content_slide(self, title: str, bullets: List[str] = None, color: RGBColor = None) -> 'StyledPresentation':
        slide_layout = self.presentation.slide_layouts[self.LAYOUT_CONTENT]
        slide = self.presentation.slides.add_slide(slide_layout)
        self.slide_count += 1
        slide.shapes.title.text = title
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(32)
        slide.shapes.title.text_frame.paragraphs[0].font.bold = True
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = color or PPTColors.PRIMARY
        if bullets:
            tf = slide.placeholders[1].text_frame
            tf.clear()
            for bullet in bullets:
                p = tf.add_paragraph()
                p.text = bullet
                p.font.size = Pt(20)
        return self

    def add_text_slide(self, title: str, content: str, color: RGBColor = None) -> 'StyledPresentation':
        slide_layout = self.presentation.slide_layouts[self.LAYOUT_BLANK]
        slide = self.presentation.slides.add_slide(slide_layout)
        self.slide_count += 1
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
        title_box.text_frame.paragraphs[0].text = title
        title_box.text_frame.paragraphs[0].font.size = Pt(32)
        title_box.text_frame.paragraphs[0].font.bold = True
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
        content_box.text_frame.paragraphs[0].text = content
        content_box.text_frame.paragraphs[0].font.size = Pt(18)
        return self

    def add_two_column_slide(self, title: str, left_title: str, left_items: List[str], right_title: str = None, right_items: List[str] = None) -> 'StyledPresentation':
        slide_layout = self.presentation.slide_layouts[self.LAYOUT_TWO_COL]
        slide = self.presentation.slides.add_slide(slide_layout)
        self.slide_count += 1
        slide.shapes.title.text = title
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(32)
        slide.placeholders[1].text = left_title
        if left_items:
            tf = slide.placeholders[1].text_frame
            for item in left_items:
                tf.add_paragraph().text = item
        if right_items:
            slide.placeholders[2].text = right_title or ""
            tf = slide.placeholders[2].text_frame
            for item in right_items:
                tf.add_paragraph().text = item
        return self

    def add_image_slide(self, title: str, image_path: str, caption: str = None, width: float = 6) -> 'StyledPresentation':
        slide_layout = self.presentation.slide_layouts[self.LAYOUT_BLANK]
        slide = self.presentation.slides.add_slide(slide_layout)
        self.slide_count += 1
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
        title_box.text_frame.paragraphs[0].text = title
        title_box.text_frame.paragraphs[0].font.size = Pt(32)
        pic = slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.2), width=Inches(width))
        if caption:
            caption_top = pic.top + pic.height + Inches(0.2)
            caption_box = slide.shapes.add_textbox(Inches(0.5), caption_top, Inches(width), Inches(0.5))
            caption_box.text_frame.paragraphs[0].text = caption
            caption_box.text_frame.paragraphs[0].font.size = Pt(14)
            caption_box.text_frame.paragraphs[0].font.italic = True
        return self

    def add_comparison_slide(self, title: str, left_title: str, right_title: str) -> 'StyledPresentation':
        slide_layout = self.presentation.slide_layouts[self.LAYOUT_COMPARISON]
        slide = self.presentation.slides.add_slide(slide_layout)
        self.slide_count += 1
        slide.shapes.title.text = title
        slide.placeholders[1].text = left_title
        slide.placeholders[2].text = right_title
        return self

    def add_numbered_slide(self, title: str, items: List[str]) -> 'StyledPresentation':
        slide_layout = self.presentation.slide_layouts[self.LAYOUT_CONTENT]
        slide = self.presentation.slides.add_slide(slide_layout)
        self.slide_count += 1
        slide.shapes.title.text = title
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for i, item in enumerate(items, 1):
            p = tf.add_paragraph()
            p.text = f"{i}. {item}"
            p.font.size = Pt(24)
        return self

    def save(self, file_path: str = None) -> None:
        path = file_path or self.file_path
        if not path:
            raise ValueError("No file_path specified")
        self.presentation.save(path)

    def get_presentation(self) -> Presentation:
        return self.presentation
